# ocr_utils.py

import pytesseract
from PIL import Image
import pdfplumber
import fitz  # pymupdf
import io

def image_to_text(image_bytes):
    """แปลงไฟล์ภาพ (bytes) เป็นข้อความด้วย OCR"""
    try:
        img = Image.open(io.BytesIO(image_bytes))
        text = pytesseract.image_to_string(img, lang='tha+eng')
        return text.strip() or "❌ ไม่พบข้อความในภาพ"
    except Exception as e:
        return f"❌ OCR image error: {e}"

def pdf_to_text(pdf_bytes):
    """แปลงไฟล์ PDF (bytes) เป็นข้อความ OCR/Extract ทั้งหมด"""
    texts = []
    try:
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    texts.append(t)
        return "\n\n".join(texts) if texts else "❌ ไม่พบข้อความใน PDF"
    except Exception as e:
        return f"❌ OCR PDF error: {e}"

def pptx_to_text(pptx_bytes):
    """ดึงข้อความจาก PowerPoint"""
    from pptx import Presentation
    texts = []
    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts) if texts else "❌ ไม่พบข้อความใน PowerPoint"
    except Exception as e:
        return f"❌ PPTX extract error: {e}"

def excel_to_text(excel_bytes):
    """ดึงข้อมูลข้อความใน Excel"""
    import pandas as pd
    try:
        xls = pd.ExcelFile(io.BytesIO(excel_bytes))
        texts = []
        for sheet in xls.sheet_names:
            df = pd.read_excel(xls, sheet)
            texts.append(f"--- Sheet: {sheet} ---\n")
            texts.append(df.to_string())
        return "\n\n".join(texts) if texts else "❌ ไม่พบข้อมูลใน Excel"
    except Exception as e:
        return f"❌ Excel extract error: {e}"
