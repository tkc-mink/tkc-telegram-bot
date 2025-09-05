# doc_extract_utils.py
# -*- coding: utf-8 -*-
from __future__ import annotations
from typing import List, Iterable, Optional
import os

# -------------------- Common helpers --------------------
def _clean_text(s: Optional[str]) -> str:
    if not s:
        return ""
    # normalize whitespace/newlines; ป้องกัน None/ซ้ำ
    out = s.replace("\r\n", "\n").replace("\r", "\n")
    # ลดช่องว่างแถวละเล็กน้อย (ไม่ aggressive เพื่อคงรูป)
    return "\n".join(line.strip() for line in out.split("\n"))

# -------------------- PDF --------------------
def extract_text_pdf(path: str, max_pages: int = 10, password: Optional[str] = None) -> str:
    """
    อ่านข้อความจาก PDF (PyPDF2)
    - รองรับเอกสารเข้ารหัส (ลอง decrypt ด้วยรหัสผ่านที่ให้ หรือว่าง)
    - ตัดที่ max_pages (ดีฟอลต์ 10 เหมือนเดิม)
    """
    try:
        from PyPDF2 import PdfReader
    except Exception as e:
        print(f"[doc_extract.pdf] PyPDF2 not available: {e}")
        return ""

    try:
        reader = PdfReader(path)
        # encrypted?
        try:
            if getattr(reader, "is_encrypted", False):
                # ลองด้วย password ที่ส่งมา และ password ว่าง
                if password:
                    try:
                        reader.decrypt(password)
                    except Exception:
                        pass
                try:
                    reader.decrypt("")  # บางไฟล์ปลดได้ด้วยค่าว่าง
                except Exception:
                    pass
        except Exception:
            pass

        texts: List[str] = []
        for i, page in enumerate(reader.pages):
            if i >= max_pages:
                texts.append("\n(ตัดเหลือ {} หน้าแรก)".format(max_pages))
                break
            try:
                t = page.extract_text() or ""
            except Exception as e:
                print(f"[doc_extract.pdf] page {i} extract error: {e}")
                t = ""
            texts.append(t)
        return _clean_text("\n".join(texts)).strip()
    except Exception as e:
        print(f"[doc_extract.pdf] read error: {e}")
        return ""

# -------------------- DOCX --------------------
def extract_text_docx(path: str, include_tables: bool = True) -> str:
    """
    อ่านข้อความจาก DOCX (python-docx)
    - include_tables=True จะดึงข้อความในตารางด้วย
    """
    try:
        from docx import Document
    except Exception as e:
        print(f"[doc_extract.docx] python-docx not available: {e}")
        return ""

    try:
        doc = Document(path)
        parts: List[str] = []
        # paragraphs
        for p in doc.paragraphs:
            parts.append(p.text or "")

        # tables (ออปชัน)
        if include_tables:
            for tb in doc.tables:
                for row in tb.rows:
                    cells = [c.text or "" for c in row.cells]
                    parts.append(" | ".join(cells))

        return _clean_text("\n".join(parts)).strip()
    except Exception as e:
        print(f"[doc_extract.docx] read error: {e}")
        return ""

# -------------------- XLSX --------------------
def extract_text_xlsx(
    path: str,
    max_rows_per_sheet: int = 500,
    max_sheets: int = 20,
) -> str:
    """
    อ่านข้อความจาก XLSX (openpyxl) — โหมด read_only
    - จำกัดแถวต่อชีทและจำนวนชีทเพื่อกันไฟล์ใหญ่
    - ใช้ data_only=True เพื่ออ่านค่าที่คำนวณแล้ว (ถ้ามี)
    """
    try:
        from openpyxl import load_workbook
    except Exception as e:
        print(f"[doc_extract.xlsx] openpyxl not available: {e}")
        return ""

    try:
        wb = load_workbook(path, read_only=True, data_only=True)
        lines: List[str] = []
        for si, ws in enumerate(wb.worksheets):
            if si >= max_sheets:
                lines.append("(ตัดจำนวนชีทเกิน {})".format(max_sheets))
                break
            lines.append(f"[Sheet] {ws.title}")
            count = 0
            for row in ws.iter_rows(values_only=True):
                # แปลงค่าให้เป็นสตริงอ่านง่าย
                line = " | ".join("" if v is None else str(v) for v in row)
                lines.append(line)
                count += 1
                if count >= max_rows_per_sheet:
                    lines.append(f"(ตัดที่ {max_rows_per_sheet} แถวแรกของชีท)")
                    break
            lines.append("----------")
        # เพื่อความเข้ากันได้กับของเดิม (เคยจำกัด 200 บรรทัดรวม)
        # ถ้าอยากเข้มเท่าเดิม เปิดบรรทัดนี้:
        # return "\n".join(lines[:200])
        return _clean_text("\n".join(lines)).strip()
    except Exception as e:
        print(f"[doc_extract.xlsx] read error: {e}")
        return ""

# -------------------- PPTX --------------------
def extract_text_pptx(path: str, include_notes: bool = True) -> str:
    """
    อ่านข้อความจาก PPTX (python-pptx)
    - ดึงข้อความจาก shapes (textbox, placeholders, tables)
    - include_notes=True จะดึง speaker notes ด้วย
    """
    try:
        from pptx import Presentation
    except Exception as e:
        print(f"[doc_extract.pptx] python-pptx not available: {e}")
        return ""

    def _shape_text(shape) -> List[str]:
        out: List[str] = []
        try:
            # text frames
            if getattr(shape, "has_text_frame", False):
                tf = shape.text_frame
                if tf:
                    for p in tf.paragraphs:
                        out.append("".join(run.text or "" for run in p.runs) or p.text or "")
            # tables
            if getattr(shape, "has_table", False):
                try:
                    tbl = shape.table
                    for r in tbl.rows:
                        cells = [(c.text or "") for c in r.cells]
                        out.append(" | ".join(cells))
                except Exception:
                    pass
            # groups (recursive)
            if hasattr(shape, "shapes"):
                for sh in shape.shapes:
                    out.extend(_shape_text(sh))
        except Exception:
            # เผื่อ shape แปลก ๆ
            pass
        return out

    try:
        prs = Presentation(path)
        lines: List[str] = []
        for idx, slide in enumerate(prs.slides, 1):
            lines.append(f"[Slide {idx}]")
            # shapes text
            for shp in slide.shapes:
                lines.extend(_shape_text(shp))
            # speaker notes
            if include_notes and getattr(slide, "has_notes_slide", False):
                try:
                    notes = slide.notes_slide.notes_text_frame
                    if notes:
                        lines.append("[Notes]")
                        for p in notes.paragraphs:
                            lines.append("".join(run.text or "" for run in p.runs) or p.text or "")
                except Exception:
                    pass
            lines.append("---")
        return _clean_text("\n".join(lines)).strip()
    except Exception as e:
        print(f"[doc_extract.pptx] read error: {e}")
        return ""

# -------------------- Dispatcher (optional helper) --------------------
def extract_text_any(path: str, **kwargs) -> str:
    """
    เดา type จากนามสกุลไฟล์แล้วเรียกฟังก์ชันที่เหมาะสม
    รองรับ: .pdf .docx .xlsx .pptx
    kwargs จะถูกส่งต่อให้ฟังก์ชันปลายทาง (เช่น max_pages, include_tables)
    """
    ext = os.path.splitext(path.lower())[1]
    if ext == ".pdf":
        return extract_text_pdf(path, **{k: v for k, v in kwargs.items() if k in {"max_pages", "password"}})
    if ext == ".docx":
        return extract_text_docx(path, **{k: v for k, v in kwargs.items() if k in {"include_tables"}})
    if ext == ".xlsx":
        return extract_text_xlsx(path, **{k: v for k, v in kwargs.items() if k in {"max_rows_per_sheet", "max_sheets"}})
    if ext == ".pptx":
        return extract_text_pptx(path, **{k: v for k, v in kwargs.items() if k in {"include_notes"}})
    print(f"[doc_extract.any] Unsupported extension: {ext}")
    return ""
