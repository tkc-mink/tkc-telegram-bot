# handlers/doc.py
# -*- coding: utf-8 -*-
"""
สรุปไฟล์ที่ผู้ใช้ส่งเข้ามาทาง Telegram (PDF / DOCX / XLSX / PPTX / TXT)
ใช้ร่วมกับ utils.message_utils.send_message และ history_utils.log_message
"""

import os
import tempfile
import requests
from typing import Optional

from PyPDF2 import PdfReader

from utils.message_utils import send_message
from history_utils import log_message
from function_calling import summarize_text_with_gpt  # ฟังก์ชันสรุปข้อความด้วย OpenAI

# ---------------------------
# Public entry point
# ---------------------------
def handle_doc(chat_id: int, msg: dict) -> None:
    """
    รับ dict ของ message จาก Telegram แล้วสรุปเอกสารที่แนบมา
    """
    doc = msg.get("document") or {}
    if not doc:
        send_message(chat_id, "❌ ไม่พบไฟล์ที่แนบมา")
        return

    file_id   = doc.get("file_id")
    file_name = doc.get("file_name", "document")
    user_id   = str(chat_id)

    # ดาวน์โหลดไฟล์ลง temp
    local_path = _download_telegram_file(file_id, file_name)
    if not local_path:
        send_message(chat_id, "❌ ไม่สามารถดาวน์โหลดไฟล์ได้")
        return

    try:
        ext = os.path.splitext(file_name)[1].lower()
        if ext == ".pdf":
            text = _extract_text_pdf(local_path)
            summary = summarize_text_with_gpt(text)

        elif ext == ".docx":
            text = _extract_text_docx(local_path)
            summary = summarize_text_with_gpt(text)

        elif ext == ".txt":
            with open(local_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            summary = summarize_text_with_gpt(text)

        elif ext == ".xlsx":
            text = _extract_text_xlsx(local_path)
            summary = summarize_text_with_gpt("ข้อมูลใน Excel:\n" + text)

        elif ext == ".pptx":
            text = _extract_text_pptx(local_path)
            summary = summarize_text_with_gpt("ข้อมูลใน PowerPoint:\n" + text)

        else:
            send_message(chat_id, "❌ รองรับเฉพาะ PDF, DOCX, XLSX, PPTX, TXT เท่านั้น")
            return

        send_message(chat_id, f"📄 สรุปไฟล์ {file_name} :\n{summary}")
        log_message(user_id, f"สรุปไฟล์ {file_name}", summary)

    except Exception as e:
        send_message(chat_id, f"❌ ไม่สามารถสรุปไฟล์ได้: {e}")
    finally:
        # ลบไฟล์ temp
        try:
            if os.path.exists(local_path):
                os.remove(local_path)
        except Exception:
            pass


# ---------------------------
# Helpers
# ---------------------------

def _download_telegram_file(file_id: str, file_name: str) -> Optional[str]:
    """
    ดึงไฟล์จาก Telegram แล้วคืน path ชั่วคราว
    """
    try:
        TELEGRAM_TOKEN = os.getenv("TELEGRAM_TOKEN")
        if not TELEGRAM_TOKEN:
            raise RuntimeError("TELEGRAM_TOKEN is not set")

        # ขอ path ของไฟล์จาก Telegram
        r = requests.get(
            f"https://api.telegram.org/bot{TELEGRAM_TOKEN}/getFile",
            params={"file_id": file_id},
            timeout=10
        )
        r.raise_for_status()
        file_path = r.json()["result"]["file_path"]

        # ดาวน์โหลดจริง
        url = f"https://api.telegram.org/file/bot{TELEGRAM_TOKEN}/{file_path}"
        suffix = os.path.splitext(file_name)[1]
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        tmp_path = tmp.name
        tmp.close()

        with requests.get(url, stream=True, timeout=20) as resp, open(tmp_path, "wb") as out:
            resp.raise_for_status()
            for chunk in resp.iter_content(chunk_size=8192):
                out.write(chunk)

        return tmp_path
    except Exception as e:
        print(f"[download_telegram_file] {e}")
        return None


def _extract_text_pdf(pdf_path: str, max_pages: int = 10) -> str:
    text = ""
    try:
        reader = PdfReader(pdf_path)
        for i, page in enumerate(reader.pages):
            if i >= max_pages:
                text += "\n(ตัดเหลือ 10 หน้าแรก)"
                break
            text += page.extract_text() or ""
    except Exception as e:
        print(f"[extract_pdf] {e}")
    return text.strip()


def _extract_text_docx(docx_path: str) -> str:
    try:
        from docx import Document
        doc = Document(docx_path)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        print(f"[extract_docx] {e}")
        return ""


def _extract_text_xlsx(xlsx_path: str) -> str:
    """
    ดึงข้อมูลเป็นข้อความจาก Excel (ตัดแค่ ~200 แถวเพื่อไม่ให้ยาวเกิน)
    """
    try:
        from openpyxl import load_workbook
        wb = load_workbook(xlsx_path, read_only=True, data_only=True)
        rows = []
        count = 0
        for ws in wb.worksheets:
            rows.append(f"=== Sheet: {ws.title} ===")
            for row in ws.iter_rows(values_only=True):
                line = " | ".join([str(c) if c is not None else "" for c in row])
                rows.append(line)
                count += 1
                if count >= 200:
                    rows.append("(ตัดความยาว Excel)")
                    break
            if count >= 200:
                break
        return "\n".join(rows)
    except Exception as e:
        print(f"[extract_xlsx] {e}")
        return ""


def _extract_text_pptx(pptx_path: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        texts = []
        for slide_idx, slide in enumerate(prs.slides, start=1):
            texts.append(f"--- Slide {slide_idx} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception as e:
        print(f"[extract_pptx] {e}")
        return ""
