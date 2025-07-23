import os
import json
import requests
import tempfile
import threading
from datetime import datetime, time as dtime
from openai import OpenAI
from PIL import Image
from PyPDF2 import PdfReader

from search_utils      import robust_image_search
from review_utils      import set_review, need_review_today
from history_utils     import log_message, get_user_history
from weather_utils     import get_weather_forecast
from gold_utils        import get_gold_price
from news_utils        import get_news
from serp_utils        import get_stock_info, get_oil_price, get_lottery_result, get_crypto_price
from function_calling  import process_with_function_calling

import backup_utils   # <-- ต้องมี backup_utils.py ตามที่จัดให้ (Google Drive)

TELEGRAM_TOKEN = os.getenv("TELEGRAM_TOKEN")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
client         = OpenAI(api_key=OPENAI_API_KEY)

USAGE_FILE           = "usage.json"
IMAGE_USAGE_FILE     = "image_usage.json"
CONTEXT_FILE         = "context_history.json"
LOCATION_FILE        = "location_logs.json"
MAX_QUESTION_PER_DAY = 30
MAX_IMAGE_PER_DAY    = 15
EXEMPT_USER_IDS      = ["6849909227"]

# --- Restore from Google Drive ทุกครั้งที่รัน ---
def try_restore_data_from_gdrive():
    try:
        print("[INFO] Attempt restore all data from Google Drive...")
        backup_utils.restore_all()
        print("[INFO] Restore complete.")
    except Exception as e:
        print(f"[ERROR] Restore fail: {e}")

try_restore_data_from_gdrive()

# --- Schedule daily backup (00:09 AM) ---
def schedule_daily_backup():
    def backup_job():
        while True:
            now = datetime.now()
            target = now.replace(hour=0, minute=9, second=0, microsecond=0)
            if now >= target:
                target = target.replace(day=now.day + 1)  # next day
                if target.month != now.month:
                    target = target.replace(month=now.month + 1, day=1)
            wait_sec = (target - now).total_seconds()
            if wait_sec <= 0:  # time already passed
                wait_sec += 86400
            print(f"[BACKUP] Waiting {wait_sec/60:.1f} minutes until next backup at {target}")
            threading.Event().wait(wait_sec)
            print("[BACKUP] Start daily backup to Google Drive...")
            backup_utils.backup_all()
            print("[BACKUP] Backup completed.")
    t = threading.Thread(target=backup_job, daemon=True)
    t.start()
schedule_daily_backup()

# --- JSON helpers ---
def load_json_safe(path):
    try:
        with open(path, "r", encoding="utf-8") as f:
            return json.load(f)
    except:
        return {}

def save_json_safe(data, path):
    try:
        with open(path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    except Exception as e:
        print(f"[save_json_safe:{path}] {e}")

# --- Usage counting ---
def check_and_increase_usage(user_id, filepath, limit):
    today = datetime.now().strftime("%Y-%m-%d")
    usage = load_json_safe(filepath)
    usage.setdefault(today, {})
    usage[today].setdefault(user_id, 0)
    if usage[today][user_id] >= limit:
        return False
    usage[today][user_id] += 1
    save_json_safe(usage, filepath)
    return True

# --- Context helpers ---
def load_context():
    return load_json_safe(CONTEXT_FILE)
def save_context(ctx):
    save_json_safe(ctx, CONTEXT_FILE)
def update_context(user_id, text):
    ctx = load_context()
    ctx.setdefault(user_id, []).append(text)
    ctx[user_id] = ctx[user_id][-6:]
    save_context(ctx)
def get_context(user_id):
    return load_context().get(user_id, [])
def is_waiting_review(user_id):
    ctx = get_context(user_id)
    return ctx and ctx[-1] == "__wait_review__"

# --- Location helpers ---
def load_location():
    return load_json_safe(LOCATION_FILE)
def save_location(loc):
    save_json_safe(loc, LOCATION_FILE)
def update_location(user_id, lat, lon):
    loc = load_location()
    loc[user_id] = {"lat": lat, "lon": lon, "ts": datetime.now().isoformat()}
    save_location(loc)
def get_user_location(user_id):
    return load_location().get(user_id)

# --- Telegram Send ---
def send_message(chat_id, text):
    try:
        requests.post(
            f"https://api.telegram.org/bot{TELEGRAM_TOKEN}/sendMessage",
            json={"chat_id": chat_id, "text": text[:4096]},
            timeout=10
        )
    except Exception as e:
        print(f"[send_message] {e}")

def send_photo(chat_id, photo_url, caption=None):
    payload = {"chat_id": chat_id, "photo": photo_url}
    if caption:
        payload["caption"] = caption
    try:
        requests.post(
            f"https://api.telegram.org/bot{TELEGRAM_TOKEN}/sendPhoto",
            json=payload,
            timeout=10
        )
    except Exception as e:
        print(f"[send_photo] {e}")

def ask_for_location(chat_id, text="📍 กรุณาแชร์ตำแหน่งของคุณ"):
    keyboard = {
        "keyboard": [
            [ {"text": "📍 แชร์ตำแหน่งของคุณ", "request_location": True} ]
        ],
        "resize_keyboard": True,
        "one_time_keyboard": True
    }
    payload = {
        "chat_id": chat_id,
        "text": text,
        "reply_markup": keyboard
    }
    try:
        requests.post(
            f"https://api.telegram.org/bot{TELEGRAM_TOKEN}/sendMessage",
            json=payload,
            timeout=5
        )
    except Exception as e:
        print(f"[ask_for_location] {e}")

def should_reset_context(new_text, prev_context):
    if not prev_context:
        return False
    last = prev_context[-1] if isinstance(prev_context, list) and prev_context else ""
    topics = ["ทอง", "หวย", "อากาศ", "ข่าว", "หุ้น", "น้ำมัน", "สุขภาพ", "ฟุตบอล"]
    if any(t in last for t in topics) and not any(t in new_text for t in topics):
        return True
    if new_text.strip().lower() in ["/reset", "เริ่มใหม่", "รีเซ็ต"]:
        return True
    return False

def handle_image_search(chat_id, user_id, text, ctx):
    if user_id not in EXEMPT_USER_IDS:
        if not check_and_increase_usage(user_id, IMAGE_USAGE_FILE, MAX_IMAGE_PER_DAY):
            send_message(chat_id, f"❌ ครบ {MAX_IMAGE_PER_DAY} รูปวันนี้แล้ว")
            return
    kw = text
    imgs = robust_image_search(kw)
    if imgs:
        for url in imgs[:3]:
            send_photo(chat_id, url, caption=f"ผลลัพธ์: {kw}")
    else:
        send_message(chat_id, f"ไม่พบภาพสำหรับ '{kw}'")

def _download_telegram_file(file_id):
    r = requests.get(
        f"https://api.telegram.org/bot{TELEGRAM_TOKEN}/getFile",
        params={"file_id": file_id},
        timeout=10
    )
    file_path = r.json()["result"]["file_path"]
    url = f"https://api.telegram.org/file/bot{TELEGRAM_TOKEN}/{file_path}"
    ext = os.path.splitext(file_path)[1].lower()
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
    with requests.get(url, stream=True, timeout=20) as resp:
        for chunk in resp.iter_content(chunk_size=8192):
            tmp.write(chunk)
    tmp.close()
    return tmp.name

def _extract_text_from_pdf(pdf_path, max_pages=10):
    try:
        reader = PdfReader(pdf_path)
        text = ""
        for i, page in enumerate(reader.pages):
            if i >= max_pages:
                text += "\n(ตัดเหลือ 10 หน้าแรก)"
                break
            text += page.extract_text() or ""
        return text.strip()
    except Exception as e:
        print(f"[PDF extract] {e}")
        return ""

def _extract_text_from_docx(docx_path):
    try:
        from docx import Document
        doc = Document(docx_path)
        return "\n".join([p.text for p in doc.paragraphs])
    except Exception as e:
        print(f"[DOCX extract] {e}")
        return ""

def _extract_text_from_xlsx(xlsx_path):
    try:
        from openpyxl import load_workbook
        wb = load_workbook(xlsx_path, read_only=True)
        text = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                line = " | ".join([str(cell) if cell is not None else "" for cell in row])
                text.append(line)
            text.append("----------")
        return "\n".join(text[:200])  # limit
    except Exception as e:
        print(f"[XLSX extract] {e}")
        return ""

def _extract_text_from_pptx(pptx_path):
    try:
        from pptx import Presentation
        prs = Presentation(pptx_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
            text.append("---")
        return "\n".join(text)
    except Exception as e:
        print(f"[PPTX extract] {e}")
        return ""

def handle_document(chat_id, doc_file, file_name, user_text=""):
    summary = ""
    ext = os.path.splitext(file_name)[1].lower()
    if ext in [".pdf"]:
        extracted = _extract_text_from_pdf(doc_file)
        if not extracted:
            send_message(chat_id, "❌ ไม่สามารถอ่าน PDF ได้")
            return
        summary = summarize_text_with_gpt(extracted)
    elif ext in [".docx"]:
        extracted = _extract_text_from_docx(doc_file)
        if not extracted:
            send_message(chat_id, "❌ ไม่สามารถอ่าน Word ได้")
            return
        summary = summarize_text_with_gpt(extracted)
    elif ext in [".txt"]:
        with open(doc_file, encoding="utf-8", errors="ignore") as f:
            text = f.read()
        summary = summarize_text_with_gpt(text)
    elif ext in [".xlsx"]:
        extracted = _extract_text_from_xlsx(doc_file)
        if not extracted:
            send_message(chat_id, "❌ ไม่สามารถอ่าน Excel ได้")
            return
        summary = summarize_text_with_gpt("ข้อมูลใน Excel:\n" + extracted)
    elif ext in [".pptx"]:
        extracted = _extract_text_from_pptx(doc_file)
        if not extracted:
            send_message(chat_id, "❌ ไม่สามารถอ่าน PowerPoint ได้")
            return
        summary = summarize_text_with_gpt("ข้อมูลใน PowerPoint:\n" + extracted)
    else:
        send_message(chat_id, "❌ รองรับ PDF, Word, Excel, PowerPoint, TXT เท่านั้น")
        return
    send_message(chat_id, f"สรุปเอกสาร/ไฟล์ {file_name} :\n{summary}")

def summarize_text_with_gpt(text, prompt="สรุปเนื้อหานี้ (ภาษาไทย):"):
    try:
        result = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "คุณคือ AI ช่วยสรุปเนื้อหาและอ่านเอกสารภาษาไทยได้ดี"},
                {"role": "user", "content": f"{prompt}\n\n{text[:5000]}"}
            ],
            max_tokens=600
        )
        return result.choices[0].message.content.strip()
    except Exception as e:
        print(f"[summarize_text_with_gpt] {e}")
        return "❌ ไม่สามารถสรุปได้"

def handle_message(data):
    msg = data.get("message", {})
    chat_id = msg.get("chat", {}).get("id")
    if not chat_id:
        return
    user_text = msg.get("caption", "") or msg.get("text", "")
    user_id   = str(chat_id)

    # 1) Location
    if "location" in msg:
        lat = msg["location"].get("latitude")
        lon = msg["location"].get("longitude")
        if lat is not None and lon is not None:
            update_location(user_id, lat, lon)
            send_message(chat_id, "✅ บันทึกตำแหน่งแล้ว! ลองถามอากาศอีกครั้งได้เลย")
        else:
            send_message(chat_id, "❌ ตำแหน่งไม่ถูกต้อง กรุณาส่งใหม่")
        return

    # 2) Document/File
    if "document" in msg:
        doc = msg["document"]
        file_id = doc["file_id"]
        file_name = doc.get("file_name", "document")
        file_path = _download_telegram_file(file_id)
        send_message(chat_id, f"🕵️‍♂️ กำลังอ่านไฟล์ {file_name} ...")
        handle_document(chat_id, file_path, file_name, user_text)
        os.remove(file_path)
        return

    # 3) Photo/Image (เฉพาะค้นหารูป ไม่รับวิเคราะห์เนื้อหา/ข้อความ)
    if "photo" in msg:
        send_message(chat_id, "ขออภัย ขณะนี้ยังไม่รองรับการอ่านข้อความหรือวิเคราะห์เนื้อหาจากรูปภาพโดยตรง หากต้องการค้นหารูปภาพให้พิมพ์ว่า 'ขอรูป...' หรือคำค้นภาพที่ต้องการ")
        return

    # /reset context manual
    if user_text.strip().lower() in ["/reset", "เริ่มใหม่", "รีเซ็ต"]:
        save_context({user_id: []})
        send_message(chat_id, "🧹 เริ่มต้นสนทนาใหม่แล้วครับ!")
        return

    if user_text.strip() == "📍 แชร์ตำแหน่งของคุณ":
        ask_for_location(chat_id)
        return

    ctx = get_context(user_id)
    if should_reset_context(user_text, ctx):
        ctx = []
        save_context({user_id: []})
    update_context(user_id, user_text)
    ctx = get_context(user_id)

    if user_text.strip() == "/my_history":
        history = get_user_history(user_id, limit=10)
        if not history:
            send_message(chat_id, "ยังไม่มีประวัติการถาม-ตอบของคุณ")
        else:
            out = "\n\n".join(f"[{it['date']}] ❓{it['q']}\n➡️ {it['a']}" for it in history)
            send_message(chat_id, f"ประวัติ 10 ล่าสุด:\n\n{out}")
        return

    if need_review_today(user_id) and not is_waiting_review(user_id):
        send_message(chat_id, "❓ กรุณารีวิววันนี้ (1-5):")
        update_context(user_id, "__wait_review__")
        return
    if is_waiting_review(user_id) and user_text.strip() in ["1","2","3","4","5"]:
        set_review(user_id, int(user_text.strip()))
        send_message(chat_id, "✅ ขอบคุณสำหรับรีวิวครับ!")
        return

    if user_id not in EXEMPT_USER_IDS:
        if not check_and_increase_usage(user_id, USAGE_FILE, MAX_QUESTION_PER_DAY):
            send_message(chat_id, f"❌ ครบ {MAX_QUESTION_PER_DAY} คำถามแล้ววันนี้")
            return

    txt = user_text.lower()
    loc = get_user_location(user_id)

    # สภาพอากาศ
    if "อากาศ" in txt or "weather" in txt:
        if loc and loc.get("lat") and loc.get("lon"):
            reply = get_weather_forecast(text=None, lat=loc["lat"], lon=loc["lon"])
            send_message(chat_id, reply)
        else:
            ask_for_location(chat_id)
        return

    # ค้นหารูป
    if any(k in txt for k in ["ขอรูป","รูป","image","photo"]):
        handle_image_search(chat_id, user_id, user_text, ctx)
        log_message(user_id, user_text, "ส่งรูปภาพ (ดูในแชท)")
        return

    # ฟังก์ชันอื่นๆ ผ่าน GPT/Function calling
    try:
        reply = process_with_function_calling(user_text, ctx=ctx[-4:])
    except Exception as e:
        print(f"[GPT function_calling] {e}")
        reply = "❌ ระบบขัดข้อง ลองใหม่อีกครั้ง"

    log_message(user_id, user_text, reply)
    send_message(chat_id, reply)
